"""
Build a fully-editable PowerPoint deck from the Brightwayz volunteer pitch.

The Marp-rendered .pptx is image-based and not editable in Google Slides.
This script produces clean editable slides using python-pptx — every text
frame is selectable and editable in PowerPoint / Google Slides / Keynote.

Run:  python build_editable_pptx.py
Output: brightwayz-volunteer-pitch-editable.pptx (alongside this script)
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

INDIGO = RGBColor(0x4F, 0x46, 0xE5)
GRAY_900 = RGBColor(0x1F, 0x29, 0x37)
GRAY_600 = RGBColor(0x6B, 0x72, 0x80)
GRAY_400 = RGBColor(0x9C, 0xA3, 0xAF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=GRAY_900, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_bulleted_list(slide, x, y, w, h, items, *, size=16, color=GRAY_900):
    """items: list of either str (bullet) or dict {text, bold, color, indent}."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, dict):
            run = p.add_run()
            run.text = item["text"]
            run.font.size = Pt(item.get("size", size))
            run.font.bold = item.get("bold", False)
            run.font.color.rgb = item.get("color", color)
            if item.get("indent"):
                p.level = item["indent"]
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(size)
            run.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def add_notes(slide, text):
    nt = slide.notes_slide.notes_text_frame
    nt.text = text


def add_kicker(slide, x, y, text):
    add_text(slide, x, y, Inches(8), Inches(0.4), text, size=12, bold=True, color=INDIGO)


def add_title(slide, x, y, w, text, size=36):
    add_text(slide, x, y, w, Inches(1.5), text, size=size, bold=True, color=INDIGO)


def add_subtitle(slide, x, y, w, text, size=20, color=GRAY_600):
    box = slide.shapes.add_textbox(x, y, w, Inches(2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return box


def add_h2(slide, x, y, text, size=22):
    return add_text(slide, x, y, Inches(6), Inches(0.5), text, size=size, bold=True, color=GRAY_900)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    BLANK = prs.slide_layouts[6]

    # ---------------------------------------------------------------------
    # Slide 1 — Title
    # ---------------------------------------------------------------------
    s = prs.slides.add_slide(BLANK)
    add_kicker(s, Inches(0.6), Inches(0.7), "Brightwayz · Demo + Volunteer Pitch")
    add_title(
        s, Inches(0.6), Inches(1.3), Inches(12),
        "Building a community-services platform — with AI as a teammate",
        size=40,
    )
    add_subtitle(
        s, Inches(0.6), Inches(3.8), Inches(12),
        "A live look at what we shipped, how we shipped it, and how you can help.",
        size=22,
    )
    add_text(
        s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.8),
        "brightwayz · connecting people with housing, food, employment, healthcare, and legal aid",
        size=14, color=GRAY_600, italic=True,
    )
    add_notes(s, """TIMING: 60 sec

OPENING (warm + confident):
"Thanks for being here. This is a 15-minute walk through a project that proves a small group of people — many of them non-coders — can ship real production software now, with AI as a teammate. I'm going to show you what we built, how, and where you fit in."

NOTES:
- Don't say "vibe coding" yet — let it land in slide 4
- Mix of technical + non-technical in the room: keep eye contact split
- If the demo URL is live, mention "you can pull this up on your phone right now"
""")

    # ---------------------------------------------------------------------
    # Slide 2 — Why
    # ---------------------------------------------------------------------
    s = prs.slides.add_slide(BLANK)
    add_title(s, Inches(0.6), Inches(0.5), Inches(12), "Why we're building this")
    add_text(
        s, Inches(0.6), Inches(1.6), Inches(12), Inches(1),
        "Community-service orgs spend hours on intake forms, follow-ups, referrals — work that's repetitive but high-stakes.",
        size=18, color=GRAY_900,
    )
    add_text(
        s, Inches(0.6), Inches(2.7), Inches(12), Inches(0.5),
        "Brightwayz cuts the friction:",
        size=18, bold=True, color=INDIGO,
    )
    add_bulleted_list(
        s, Inches(0.8), Inches(3.4), Inches(12), Inches(3),
        [
            "• People in need talk to an AI assistant (English / Spanish / French / Chinese / Arabic), describe their situation, get matched.",
            "• Case managers see every request in one dashboard: filter by need, status, ZIP, date.",
            "• Status changes auto-notify the client — email + SMS + WhatsApp — so nobody's left wondering.",
        ],
        size=16,
    )
    add_text(
        s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
        "A free, accessible front door to social services. Built for and with the orgs that use it.",
        size=15, italic=True, color=GRAY_600,
    )
    add_notes(s, """TIMING: 90 sec

KEY POINTS:
1. Frame the pain first — "intake hell" resonates with anyone who's ever volunteered at a food bank, free clinic, or housing nonprofit.
2. Three concrete benefits, one per audience role.
3. Land on "built for and with" — this signals that we want their input, which is the recruiting hook.

ANALOGY (if non-tech audience):
"Think of it like the difference between filling out 8 PDFs and texting a friend who already knows what to ask."

DO NOT:
- List every feature here. That's the next slide.
- Mention competitors. Stay focused on the work.
""")

    # ---------------------------------------------------------------------
    # Slide 3 — What we've shipped
    # ---------------------------------------------------------------------
    s = prs.slides.add_slide(BLANK)
    add_title(s, Inches(0.6), Inches(0.5), Inches(12), "What we've shipped (so far)")
    add_h2(s, Inches(0.6), Inches(1.8), "For people seeking help")
    add_bulleted_list(
        s, Inches(0.6), Inches(2.4), Inches(6.2), Inches(4.5),
        [
            "• AI chatbot on the homepage",
            "• Multi-step intake form",
            "• Public resource directory",
            "• Referral acceptance flow (token-based)",
            "• Mobile-first, multi-language",
        ],
        size=15,
    )
    add_h2(s, Inches(7.0), Inches(1.8), "For case managers")
    add_bulleted_list(
        s, Inches(7.0), Inches(2.4), Inches(6.0), Inches(4.5),
        [
            "• Dashboard: clients · intakes · referrals · cases",
            "• Search & filter (name, ZIP, date, status)",
            "• CRM lifecycle: open → assigned → in-progress → resolved → closed",
            "• Status changes → auto email + SMS to client",
            "• One-click WhatsApp / SMS to any client",
            "• Rich-text internal notes (case-worker only)",
            "• Convert intake → referral with one click",
        ],
        size=15,
    )
    add_notes(s, """TIMING: 2 min (this is where the live demo happens)

LIVE DEMO:
Open two browser windows on the projector:
1. Phone or laptop on the public landing page (the AI chatbot)
2. Dashboard logged in as a case worker

WALK THROUGH:
- "I'm a person needing housing help" → say one sentence to the AI chatbot. Watch it ask for the right follow-up info.
- Switch to dashboard → "and here it is, in real time."
- Click into the new intake → show the case status dropdown → flip it to 'assigned' → "watch — the client just got an email saying a case manager is on it. And an SMS, and a WhatsApp."

IF DEMO BREAKS:
- That's fine — say "this is software, software breaks."
- Pivot to screenshots if you've got them prepped.
- Don't try to debug live; it eats time and audience attention.

IF DEMO GOES LONG:
- Cut the SMS / WhatsApp demo (the email one alone makes the point).
""")

    # ---------------------------------------------------------------------
    # Slide 4 — Vibe coding
    # ---------------------------------------------------------------------
    s = prs.slides.add_slide(BLANK)
    add_title(
        s, Inches(0.6), Inches(0.5), Inches(12),
        "How we built it: \"vibe coding\" with Claude Code",
        size=30,
    )
    add_text(
        s, Inches(0.6), Inches(1.8), Inches(12), Inches(1),
        "Vibe coding = describe the outcome, let AI do the typing. You review, course-correct, and own the direction.",
        size=18, italic=True, color=INDIGO,
    )
    add_h2(s, Inches(0.6), Inches(3.2), "A real exchange from this project:", size=18)
    add_bulleted_list(
        s, Inches(0.6), Inches(3.9), Inches(12), Inches(2.5),
        [
            {"text": "Human:  add a whatsapp button same as text message", "size": 16, "bold": True, "color": GRAY_900},
            {"text": "Claude Code:  (reads existing SMS code, adds Twilio WhatsApp wrapper, new endpoint, frontend button + panel, runs tests, opens PR)", "size": 16, "color": GRAY_600, "italic": True},
        ],
        size=16,
    )
    add_text(
        s, Inches(0.6), Inches(5.7), Inches(12), Inches(1),
        "Ten minutes later: green button on the dashboard, real WhatsApp delivered to a phone.",
        size=16, bold=True, color=INDIGO,
    )
    add_text(
        s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
        "What changed: the bottleneck stopped being typing code. It became knowing what to ask for — design, UX, prioritization.",
        size=13, color=GRAY_600,
    )
    add_notes(s, """TIMING: 2 min

DEFINE FIRST: pause and say the term out loud.
"Vibe coding is a term coined by Andrej Karpathy earlier this year. The idea is: you describe the vibe of what you want, and the AI handles the typing. Your job becomes direction, not production."

THEN: the WhatsApp story.
- Six words → working feature, end to end.
- Mention this is REAL — not staged, not a demo project. Real Twilio account, real WhatsApp message, real production deploy at the time.

PUNCHLINE — say it slowly:
"The bottleneck stopped being typing code. It became knowing what to ask for."

TIE TO AUDIENCE:
- Tech folks: "Your code-review skills matter more than ever."
- Non-tech: "If you can describe what you want clearly, you can ship."

OPTIONAL LIVE BIT:
If you have a laptop with Claude Code running, open a new terminal and run a tiny request live. Example: ask it to add a confirmation modal to one button. Watch it write the code in front of the audience. RISKY but high impact if it works.
""")

    # ---------------------------------------------------------------------
    # Slide 5 — Agentic AI
    # ---------------------------------------------------------------------
    s = prs.slides.add_slide(BLANK)
    add_title(s, Inches(0.6), Inches(0.5), Inches(12), "Agentic AI — beyond chat")
    add_text(
        s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.6),
        "Most AI today answers. Agentic AI acts.",
        size=22, bold=True, color=INDIGO,
    )
    add_text(
        s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.5),
        "In this project, the AI:",
        size=16, color=GRAY_900,
    )
    add_bulleted_list(
        s, Inches(0.8), Inches(3.1), Inches(12), Inches(4),
        [
            "✓ Read existing code to match conventions",
            "✓ Edited multiple files in one go",
            "✓ Ran npm, pip, pytest, git, aws, curl",
            "✓ Deployed to AWS (ECS, ALB, CloudFront, Secrets Manager)",
            "✓ Diagnosed real bugs (JWT signing algorithm mismatch, DNS misconfig)",
            "✓ Rotated leaked credentials",
            "✓ Wrote SQL migrations, applied them, verified data",
        ],
        size=15,
    )
    add_text(
        s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
        "A single human + an agent → a small engineering team. Without the meetings.",
        size=15, italic=True, color=INDIGO,
    )
    add_notes(s, """TIMING: 2 min

CONTRAST CLEARLY:
"ChatGPT answers questions. An agent does the work."

POINT TO THE LIST:
"Every check mark here is a real thing that happened on this project. Not in a sandbox — on actual AWS infrastructure with real credentials."

WAR STORY (pick one):
1. JWT mismatch — Supabase migrated to ES256 partway through, every login broke, agent diagnosed and fixed it in 10 minutes.
2. Leaked AWS keys — agent spotted them in a config file, rotated them, and updated GitHub secrets and the live ECS task — without downtime.
3. DNS — agent traced why www.wedkai.com worked but apex didn't, all the way to GoDaddy forwarding settings.

These are the kinds of multi-step incidents that used to need an on-call engineer + a meeting. Now: one prompt, one human watching.

DON'T:
- Oversell. The human still needed to confirm sensitive steps.
- The point is "force multiplier," not "replaces humans."
""")

    # ---------------------------------------------------------------------
    # Slide 6 — Under the hood
    # ---------------------------------------------------------------------
    s = prs.slides.add_slide(BLANK)
    add_title(s, Inches(0.6), Inches(0.5), Inches(12), "What's under the hood")
    add_h2(s, Inches(0.6), Inches(1.7), "Web", size=18)
    add_bulleted_list(
        s, Inches(0.6), Inches(2.2), Inches(6.2), Inches(2),
        [
            "• Next.js 16 (App Router, static export)",
            "• Tailwind CSS · Tiptap rich text",
            "• Supabase Auth (JWT ES256)",
        ],
        size=14,
    )
    add_h2(s, Inches(0.6), Inches(4.2), "API", size=18)
    add_bulleted_list(
        s, Inches(0.6), Inches(4.7), Inches(6.2), Inches(2.5),
        [
            "• FastAPI + Pydantic",
            "• Supabase Postgres",
            "• Anthropic Claude (intake assistant)",
            "• Twilio (SMS / WhatsApp)",
            "• Resend (email)",
        ],
        size=14,
    )
    add_h2(s, Inches(7.0), Inches(1.7), "Infra (when live)", size=18)
    add_bulleted_list(
        s, Inches(7.0), Inches(2.2), Inches(6.0), Inches(2),
        [
            "• AWS ECS Fargate · ALB",
            "• CloudFront + S3 (static web)",
            "• Secrets Manager",
            "• GitHub Actions CI/CD",
        ],
        size=14,
    )
    add_h2(s, Inches(7.0), Inches(4.2), "Dev", size=18)
    add_bulleted_list(
        s, Inches(7.0), Inches(4.7), Inches(6.0), Inches(2.5),
        [
            "• Local-first: npm run dev + uvicorn",
            "• venv for isolation",
            "• Branches per feature, easy to revert",
        ],
        size=14,
    )
    add_text(
        s, Inches(0.6), Inches(6.9), Inches(12), Inches(0.4),
        "Every integration above was wired up by Claude Code, including the gnarly bits (DNS, IAM, secrets rotation).",
        size=12, italic=True, color=GRAY_600,
    )
    add_notes(s, """TIMING: 60 sec — this is a "skim" slide

ASK FIRST: "Quick show of hands — who's a developer?"

IF MOSTLY DEVS:
Spend more time. Mention you can swap any layer (e.g. Resend → SES, Supabase → self-hosted Postgres) and the AI handles the refactor.

IF MOSTLY NON-DEVS:
"You don't need to know any of these names. I'm showing you the stack just so you know there's nothing exotic here — these are standard tools. Anyone with a few months of practice could maintain this."

PUNCH:
"What's interesting isn't the stack. It's that we shipped this whole stack with the AI as the integration engineer."
""")

    # ---------------------------------------------------------------------
    # Slide 7 — How you can help (non-technical)
    # ---------------------------------------------------------------------
    s = prs.slides.add_slide(BLANK)
    add_title(s, Inches(0.6), Inches(0.5), Inches(12), "How you can help — non-technical")
    add_text(
        s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5),
        "You don't need to write code to make this better.",
        size=18, bold=True, color=INDIGO,
    )
    add_bulleted_list(
        s, Inches(0.6), Inches(2.5), Inches(12), Inches(4.5),
        [
            "• Try it. Submit a fake intake. Tell us what was confusing.",
            "• Translate. Are the AI's prompts kind / clear in your language? Help us tune them.",
            "• Pilot it. Run an intake or two for your org and tell us what's missing.",
            "• Design. The landing page, the email templates, the SMS phrasing — these need a human touch.",
            "• Write. Resource directory entries, FAQ, how-to guides for case managers.",
            "• Talk to people. Spread the word to orgs and to people who need help.",
        ],
        size=15,
    )
    add_text(
        s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
        "The features above were prioritized by someone like you telling the team what mattered.",
        size=14, italic=True, color=INDIGO,
    )
    add_notes(s, """TIMING: 90 sec

POSTURE: lean in. This is the recruiting moment for half your audience.

EMPHASIZE:
"You don't need to write code." Repeat it. Many people in the room will be quietly thinking they're not technical enough to contribute.

CONCRETE ASKS:
Pick ONE bullet and go deep, rather than skimming all six. The "translate" bullet is great if the room is diverse linguistically. The "pilot it" bullet is great if there are nonprofit folks present.

CALL OUT NAMES IF YOU CAN:
"Sara, you mentioned you run intake at [Org] — that 'pilot it' bullet is literally aimed at you."

THE LAST LINE IS THE PUNCHLINE:
"The features above were prioritized by someone like you telling the team what mattered." Pause. Let it sit.
""")

    # ---------------------------------------------------------------------
    # Slide 8 — How you can help (technical)
    # ---------------------------------------------------------------------
    s = prs.slides.add_slide(BLANK)
    add_title(s, Inches(0.6), Inches(0.5), Inches(12), "How you can help — technical")
    add_text(
        s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.6),
        "If you code (or want to learn), Claude Code makes this exceptionally low-friction to contribute.",
        size=16, color=GRAY_900,
    )
    add_bulleted_list(
        s, Inches(0.6), Inches(2.6), Inches(12), Inches(4),
        [
            {"text": "• Pick a feature from the roadmap below, branch it, ship it.", "size": 15, "color": GRAY_900},
            {"text": "    – File uploads (S3 task-role auth, currently stubbed)", "size": 14, "color": GRAY_600},
            {"text": "    – Audit log view (table already collects entries)", "size": 14, "color": GRAY_600},
            {"text": "    – Dashboard analytics (intake counts, response times)", "size": 14, "color": GRAY_600},
            {"text": "    – SMS opt-out + STOP-keyword handling", "size": 14, "color": GRAY_600},
            {"text": "    – Per-org branding", "size": 14, "color": GRAY_600},
            {"text": "• Bring your own AI. Claude Code, Cursor, whatever — the codebase reads cleanly for any agent.", "size": 15, "color": GRAY_900},
            {"text": "• No prior repo knowledge needed. Open the repo in your editor, type a request to your AI: \"Add an X to Y.\" Read what it writes. Send a PR.", "size": 15, "color": GRAY_900},
        ],
        size=15,
    )
    add_notes(s, """TIMING: 90 sec

POSTURE: this is the recruiting moment for the OTHER half of the audience.

THE 5 ROADMAP ITEMS:
Don't read them all out. Mention 2-3 you find most interesting. The bullets are there for the slide deck PDF that lives on after.

REASSURE NEW CONTRIBUTORS:
"If you've never contributed to an open-source project, this is a gentle entry point. Pick a small thing. Type one sentence to your AI. Read what it produces. Send a PR. We'll review it kindly."

REASSURE EXPERIENCED ENGINEERS:
"If you're a senior eng wondering 'is this codebase a mess' — the quality bar in here is set by Pydantic + TypeScript types + pytest + ESLint. You'll feel at home."

OFFER PAIR PROGRAMMING:
"If you want to ship something together, I'll pair with you live for an hour. Just DM me."
""")

    # ---------------------------------------------------------------------
    # Slide 9 — What this demo proved
    # ---------------------------------------------------------------------
    s = prs.slides.add_slide(BLANK)
    add_title(s, Inches(0.6), Inches(0.5), Inches(12), "What this demo just proved")
    add_h2(s, Inches(0.6), Inches(1.8), "Most of what you saw was built in conversations.", size=18)
    add_bulleted_list(
        s, Inches(0.6), Inches(2.6), Inches(6.2), Inches(4),
        [
            "• 4 feature branches, all merged",
            "• 20+ committed features",
            "• 3 messaging providers integrated end-to-end",
            "• 1 full AWS deploy + teardown",
            "• A migration, a bug fix, and a UX iteration in the time it took to drink a coffee",
        ],
        size=15,
    )
    add_h2(s, Inches(7.0), Inches(1.8), "The hard parts shifted:", size=18)
    add_bulleted_list(
        s, Inches(7.0), Inches(2.6), Inches(6.0), Inches(3),
        [
            "• Typing boilerplate  →  asking the right question",
            "• \"How do I configure Twilio\"  →  \"what should the trial-vs-prod path look like\"",
            "• Spinning up infra  →  deciding what's worth deploying",
        ],
        size=15,
    )
    add_text(
        s, Inches(7.0), Inches(5.5), Inches(6), Inches(1.5),
        "The remaining hard parts are judgment, taste, and care — exactly the things humans should own.",
        size=15, italic=True, bold=True, color=INDIGO,
    )
    add_notes(s, """TIMING: 60 sec

POSTURE: this is the takeaway slide. Slow down.

NUMBERS LAND BETTER WHEN SPOKEN, NOT READ:
Don't say "twenty plus committed features." Say "we shipped more than twenty features. And that's just what's on this branch."

THE BEFORE/AFTER LIST:
Read each item aloud. The "before → after" cadence is the rhetorical trick that makes the point stick.

LAST LINE — say slowly, look up from notes:
"The remaining hard parts are judgment, taste, and care — exactly the things humans should own."

This is your money line. Pause. Move on.
""")

    # ---------------------------------------------------------------------
    # Slide 10 — Get involved
    # ---------------------------------------------------------------------
    s = prs.slides.add_slide(BLANK)
    add_title(s, Inches(0.6), Inches(0.5), Inches(12), "Get involved")
    add_h2(s, Inches(0.6), Inches(1.7), "Today", size=20)
    add_bulleted_list(
        s, Inches(0.6), Inches(2.3), Inches(6.5), Inches(2.5),
        [
            "• Try the demo · INSERT URL",
            "• Star / fork the repo · github.com/aerrami/brightwayz-{api,web}",
            "• Join our Slack / Discord · INSERT LINK",
        ],
        size=15,
    )
    add_h2(s, Inches(0.6), Inches(4.5), "This week", size=20)
    add_bulleted_list(
        s, Inches(0.6), Inches(5.1), Inches(6.5), Inches(2.5),
        [
            "• Pick one item from the \"How you can help\" lists",
            "• DM us with your name + what you'd like to try",
            "• We'll pair you with whoever's working on it",
        ],
        size=15,
    )
    add_h2(s, Inches(7.5), Inches(1.7), "Contact", size=20)
    add_bulleted_list(
        s, Inches(7.5), Inches(2.3), Inches(5.5), Inches(3),
        [
            "✉  INSERT EMAIL",
            "🌐  INSERT SITE",
            "🐙  GITHUB ORG / HANDLE",
        ],
        size=15,
    )
    add_text(
        s, Inches(7.5), Inches(5.0), Inches(5.5), Inches(0.6),
        "Thank you 💛",
        size=26, bold=True, color=INDIGO,
    )
    add_text(
        s, Inches(7.5), Inches(5.8), Inches(5.5), Inches(1),
        "Building public-good software, faster — together.",
        size=16, italic=True, color=GRAY_600,
    )
    add_notes(s, """TIMING: 90 sec + Q&A

LEAVE THIS SLIDE UP through Q&A — people will photograph it.

OFFER TO LINGER:
"I'll be in the [room / corner / coffee area] for the next 20 minutes. Come find me if you want to chat about anything you saw today."

ASK A QUESTION BACK:
"Before we close — what's one feature you'd add tomorrow if you had the time and tools? Just shout it out."
This often surfaces your next 2-3 best contributors.

PRACTICAL NEXT STEPS:
- If you're collecting signups, have a QR code or Google Form ready.
- Print 5-10 paper copies of the contact info for people who don't photograph slides.
- Mention any in-person follow-up: "I'm hosting an intro-to-Claude-Code workshop next Saturday."
""")

    out = Path(__file__).parent / "brightwayz-volunteer-pitch-editable.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
